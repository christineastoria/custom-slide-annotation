"""
Financial Slide Deck Generator Agent

A simple agent that generates polished financial slide decks as PPTX files.
The agent has granular control over slide design using individual tools.
"""

from dotenv import load_dotenv
load_dotenv(override=True)

from io import BytesIO
from typing import Optional
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from langchain.chat_models import init_chat_model
from langchain_core.tools import tool
from langchain_core.messages import HumanMessage, ToolMessage
from langchain.agents import create_agent
from langsmith import uuid7


# Initialize the model
model = init_chat_model("openai:gpt-4o-mini")


# ============================================================================
# PRESENTATION STATE (shared across tool calls)
# ============================================================================

class PresentationBuilder:
    """Manages the presentation state across tool calls."""

    def __init__(self):
        self.prs: Optional[Presentation] = None
        self.current_slide = None
        self.slide_count = 0

    def reset(self):
        """Reset the builder for a new presentation."""
        if self.prs is not None:
            del self.prs
        self.prs = None
        self.current_slide = None
        self.slide_count = 0

    def create_new(self):
        """Create a fresh presentation."""
        self.reset()
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)  # 16:9
        self.prs.slide_height = Inches(7.5)

    def get_slide(self, slide_num: int):
        """1-indexed slide access."""
        if self.prs is None:
            return None
        if slide_num < 1 or slide_num > len(self.prs.slides):
            return None
        return self.prs.slides[slide_num - 1]

    def ensure_current_slide(self, slide_num: Optional[int] = None):
        """
        If slide_num is provided, set current_slide to that slide (1-indexed).
        Otherwise, keep current_slide as-is.
        """
        if self.prs is None:
            return False
        if slide_num is None:
            return self.current_slide is not None
        slide = self.get_slide(slide_num)
        if slide is None:
            return False
        self.current_slide = slide
        return True


# Global builder instance
builder = PresentationBuilder()


# ============================================================================
# PPTX TOOLS - Individual design tools for the agent
# ============================================================================

def _parse_hex_color(hex_color: str) -> RGBColor:
    hex_color = hex_color.lstrip("#")
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    return RGBColor(r, g, b)


@tool
def create_presentation(title: str) -> str:
    """
    Create a new PowerPoint presentation. Must be called first before adding slides.
    """
    if builder.prs is not None:
        print("DEBUG: create_presentation called but presentation already exists")
        return "Presentation already exists. Continue adding slides."

    builder.create_new()
    print(f"DEBUG: Created new presentation: {title}")
    return f"Created new presentation: '{title}'"


@tool
def add_slide(background_color: str = "#0f172a") -> str:
    """
    Add a new blank slide to the presentation. Sets it as the current slide.
    """
    if builder.prs is None:
        return "Error: No presentation created. Call create_presentation first."

    blank_layout = builder.prs.slide_layouts[6]  # Blank slide
    builder.current_slide = builder.prs.slides.add_slide(blank_layout)
    builder.slide_count += 1

    print(
        f"DEBUG: add_slide called - builder.slide_count={builder.slide_count}, "
        f"prs.slides={len(builder.prs.slides)}"
    )

    # Set background
    background = builder.current_slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = _parse_hex_color(background_color)

    return f"Added slide {builder.slide_count} with background {background_color}"


@tool
def set_current_slide(slide_num: int) -> str:
    """
    Select which slide subsequent elements are added to (1-indexed).
    """
    if builder.prs is None:
        return "Error: No presentation created. Call create_presentation first."
    slide = builder.get_slide(slide_num)
    if slide is None:
        return f"Error: slide_num {slide_num} out of range. Have {len(builder.prs.slides)} slides."
    builder.current_slide = slide
    return f"Current slide set to {slide_num}"


@tool
def add_title_text(
    text: str,
    x_inches: float = 0.5,
    y_inches: float = 2.5,
    width_inches: float = 12.333,
    font_size: int = 48,
    font_color: str = "#f8fafc",
    bold: bool = True,
    center: bool = True,
    slide_num: Optional[int] = None,
) -> str:
    """
    Add a large title text to the selected/current slide.
    """
    if not builder.ensure_current_slide(slide_num):
        return "Error: No slide available. Call add_slide first (or set_current_slide)."

    textbox = builder.current_slide.shapes.add_textbox(
        Inches(x_inches), Inches(y_inches), Inches(width_inches), Inches(1.5)
    )
    tf = textbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = _parse_hex_color(font_color)
    p.font.bold = bold
    p.alignment = PP_ALIGN.CENTER if center else PP_ALIGN.LEFT

    return f"Added title: '{text}'"


@tool
def add_body_text(
    text: str,
    x_inches: float = 0.5,
    y_inches: float = 1.5,
    width_inches: float = 12.0,
    font_size: int = 16,
    font_color: str = "#94a3b8",
    slide_num: Optional[int] = None,
) -> str:
    """
    Add body/paragraph text to the selected/current slide.
    """
    if not builder.ensure_current_slide(slide_num):
        return "Error: No slide available. Call add_slide first (or set_current_slide)."

    textbox = builder.current_slide.shapes.add_textbox(
        Inches(x_inches), Inches(y_inches), Inches(width_inches), Inches(1.0)
    )
    tf = textbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = _parse_hex_color(font_color)
    p.alignment = PP_ALIGN.LEFT

    return f"Added body text: '{text[:50]}...'" if len(text) > 50 else f"Added body text: '{text}'"


@tool
def add_slide_number(
    x_inches: float = 0.5,
    y_inches: float = 0.4,
    slide_num: Optional[int] = None,
) -> str:
    """
    Add a slide number label to the selected/current slide.
    If slide_num is provided, it targets that slide; the label uses that slide's number.
    """
    if builder.prs is None:
        return "Error: No presentation created. Call create_presentation first."
    if slide_num is not None:
        if not builder.ensure_current_slide(slide_num):
            return "Error: slide_num out of range."
        label_num = slide_num
    else:
        if builder.current_slide is None:
            return "Error: No slide available. Call add_slide first."
        # derive slide index from prs
        label_num = list(builder.prs.slides).index(builder.current_slide) + 1

    textbox = builder.current_slide.shapes.add_textbox(
        Inches(x_inches), Inches(y_inches), Inches(2), Inches(0.3)
    )
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = f"Slide {label_num}"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(100, 116, 139)

    return f"Added slide number: {label_num}"


@tool
def add_metric_card(
    label: str,
    value: str,
    trend: str,
    x_inches: float,
    y_inches: float,
    width_inches: float = 2.8,
    height_inches: float = 1.3,
    slide_num: Optional[int] = None,
) -> str:
    """
    Add a styled metric card showing a KPI with trend indicator to the selected/current slide.
    """
    if not builder.ensure_current_slide(slide_num):
        return "Error: No slide available. Call add_slide first (or set_current_slide)."

    slide = builder.current_slide

    # Card background
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(x_inches), Inches(y_inches),
        Inches(width_inches), Inches(height_inches),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(30, 41, 59)  # #1e293b
    shape.line.fill.background()  # no border

    # Trend indicator
    if trend == "up":
        trend_icon = "↑"
        trend_rgb = (16, 185, 129)
    elif trend == "down":
        trend_icon = "↓"
        trend_rgb = (239, 68, 68)
    else:
        trend_icon = "→"
        trend_rgb = (107, 114, 128)

    # Label
    label_box = slide.shapes.add_textbox(
        Inches(x_inches + 0.1), Inches(y_inches + 0.15),
        Inches(width_inches - 0.2), Inches(0.3),
    )
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = label.upper()
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(148, 163, 184)
    p.alignment = PP_ALIGN.CENTER

    # Value
    value_box = slide.shapes.add_textbox(
        Inches(x_inches + 0.1), Inches(y_inches + 0.4),
        Inches(width_inches - 0.2), Inches(0.4),
    )
    tf = value_box.text_frame
    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(248, 250, 252)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Trend
    trend_box = slide.shapes.add_textbox(
        Inches(x_inches + 0.1), Inches(y_inches + 0.85),
        Inches(width_inches - 0.2), Inches(0.3),
    )
    tf = trend_box.text_frame
    p = tf.paragraphs[0]
    p.text = trend_icon
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(*trend_rgb)
    p.alignment = PP_ALIGN.CENTER

    return f"Added metric card: {label} = {value} ({trend})"


@tool
def add_subtitle(
    text: str,
    x_inches: float = 0.5,
    y_inches: float = 4.2,
    font_color: str = "#64748b",
    slide_num: Optional[int] = None,
) -> str:
    """
    Add a subtle subtitle/footer text to the selected/current slide.
    """
    if not builder.ensure_current_slide(slide_num):
        return "Error: No slide available. Call add_slide first (or set_current_slide)."

    textbox = builder.current_slide.shapes.add_textbox(
        Inches(x_inches), Inches(y_inches), Inches(12.333), Inches(0.5)
    )
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.color.rgb = _parse_hex_color(font_color)
    p.alignment = PP_ALIGN.CENTER

    return f"Added subtitle: '{text}'"


@tool(return_direct=True)
def finalize_presentation() -> bytes:
    """
    Finalize and return the completed PowerPoint presentation.
    """
    if builder.prs is None:
        return b"Error: No presentation to finalize."
    if builder.slide_count == 0:
        return b"Error: No slides added to presentation."

    print(
        f"DEBUG: finalize_presentation - builder.slide_count={builder.slide_count}, "
        f"prs.slides={len(builder.prs.slides)}"
    )

    buffer = BytesIO()
    builder.prs.save(buffer)
    buffer.seek(0)
    pptx_bytes = buffer.getvalue()

    print(f"DEBUG: finalize_presentation - saved {len(pptx_bytes)} bytes")

    builder.reset()
    return pptx_bytes


# ============================================================================
# ALL TOOLS
# ============================================================================

pptx_tools = [
    create_presentation,
    add_slide,
    set_current_slide,   # NEW
    add_title_text,
    add_body_text,
    add_slide_number,
    add_metric_card,
    add_subtitle,
    finalize_presentation,
]


# ============================================================================
# AGENT SETUP
# ============================================================================

slide_agent_prompt = """
You are a financial analyst that creates simple, focused PowerPoint decks.

TOOLS:
- create_presentation(title) - Start new deck (call FIRST)
- add_slide(background_color) - Add slide (use "#0f172a" dark blue)
- set_current_slide(slide_num) - ALWAYS select the slide you are about to edit (1-indexed)
- add_title_text(..., slide_num=) - Large heading (optionally specify slide_num)
- add_body_text(..., slide_num=) - Paragraph text
- add_metric_card(..., slide_num=) - KPI card with up/down/flat trend
- add_subtitle(..., slide_num=) - Small footer text
- add_slide_number(..., slide_num=) - Slide label
- finalize_presentation() - Save and return PPTX (call LAST)

CRITICAL RULE:
Before adding ANY content, you MUST target the correct slide:
- Either call set_current_slide(N) first, OR pass slide_num=N into the add_* tool.
Never assume the "current slide" is correct.

RULES:
1. Create EXACTLY 2-3 slides total
2. Slide 1: Title slide only (title at y=2.5, subtitle at y=4.2)
3. Slide 2-3: Content slides with 3-4 metric cards each
4. Each slide shows DIFFERENT metrics - NO duplicate data across slides
5. Use positions: title at y=1.0, metrics at y=3.0

METRIC CARD POSITIONS (for 3-4 cards in a row):
- 3 cards: x = 2.0, 5.5, 9.0
- 4 cards: x = 1.5, 4.5, 7.5, 10.5

FORMAT VALUES:
- Currency: $3.4M not $3400000
- Percentages: 24.6% with one decimal
- Trends: "up" (green), "down" (red), "flat" (gray)

Keep it simple. 2-3 slides max. No overlapping content.
"""

slide_agent = create_agent(
    model=model,
    tools=pptx_tools,
    name="financial_slide_agent",
    system_prompt=slide_agent_prompt,
).with_config({"recursion_limit": 50})


# ============================================================================
# SAMPLE DATASETS
# ============================================================================

company_financials = pd.DataFrame({
    "Quarter": ["Q1 2024", "Q2 2024", "Q3 2024", "Q4 2024"],
    "Revenue": [2500000, 2750000, 3100000, 3450000],
    "COGS": [1500000, 1600000, 1750000, 1850000],
    "Gross Profit": [1000000, 1150000, 1350000, 1600000],
    "Operating Expenses": [600000, 650000, 700000, 750000],
    "Net Income": [400000, 500000, 650000, 850000],
    "Gross Margin %": [40.0, 41.8, 43.5, 46.4],
    "Net Margin %": [16.0, 18.2, 21.0, 24.6],
    "Customers": [1200, 1350, 1520, 1750],
    "Employees": [45, 52, 58, 65]
})

saas_metrics = pd.DataFrame({
    "Month": ["Sep 2024", "Oct 2024", "Nov 2024", "Dec 2024"],
    "MRR": [95000, 112000, 135000, 165000],
    "ARR": [1140000, 1344000, 1620000, 1980000],
    "New Customers": [45, 62, 78, 95],
    "Churned Customers": [8, 6, 7, 5],
    "Net New MRR": [12000, 17000, 23000, 30000],
    "Churn Rate %": [3.5, 2.8, 2.4, 1.9],
    "CAC": [520, 485, 450, 410],
    "LTV": [4800, 5200, 5800, 6500],
    "LTV/CAC Ratio": [9.2, 10.7, 12.9, 15.9]
})

ecommerce_data = pd.DataFrame({
    "Category": ["Electronics", "Apparel", "Home & Garden", "Sports", "Beauty"],
    "Revenue": [1250000, 890000, 650000, 420000, 380000],
    "Orders": [8500, 12000, 5200, 3800, 6500],
    "Avg Order Value": [147, 74, 125, 111, 58],
    "Return Rate %": [8.5, 15.2, 4.3, 6.8, 3.2],
    "Profit Margin %": [18.5, 42.0, 35.0, 28.0, 55.0],
    "YoY Growth %": [12.0, 25.0, 8.0, 35.0, 45.0]
})


# ============================================================================
# GENERATE SLIDE DECKS
# ============================================================================

def generate_deck(data: pd.DataFrame, prompt: str) -> Optional[bytes]:
    """
    Generate a slide deck from data and return the PPTX bytes.
    """
    builder.reset()
    config = {"configurable": {"thread_id": uuid7()}}

    result = slide_agent.invoke(
        {"messages": [HumanMessage(content=f"{prompt}\n\nData:\n{data.to_string()}")]},
        config=config
    )

    for message in result["messages"]:
        if isinstance(message, ToolMessage) and message.name == "finalize_presentation":
            content = message.content
            print(f"DEBUG generate_deck: content type = {type(content)}")
            if isinstance(content, bytes):
                return content
            elif isinstance(content, str) and content.startswith("b'"):
                import ast
                try:
                    return ast.literal_eval(content)
                except Exception:
                    return None

    return None


def main():
    print("Generating slide decks...\n")

    # Deck 1
    print("1. Generating Q4 2024 Quarterly Business Review...")
    pptx_bytes = generate_deck(
        company_financials,
        """Create a Q4 2024 Business Review with exactly 3 slides:

Slide 1: Title slide - "Q4 2024 Business Review"
Slide 2: Revenue metrics (Revenue, Gross Profit, Net Income)
Slide 3: Growth metrics (Customers, Employees, Net Margin %)

Use Q4 2024 data only. Show trends vs Q3.
"""
    )
    if pptx_bytes:
        with open("Q4_2024_Business_Review.pptx", "wb") as f:
            f.write(pptx_bytes)
        print("   Saved to Q4_2024_Business_Review.pptx")
    else:
        print("   Failed to generate")

    # Deck 2
    print("2. Generating SaaS Investor Update...")
    pptx_bytes = generate_deck(
        saas_metrics,
        """Create a SaaS Investor Update with exactly 3 slides:

Slide 1: Title slide - "SaaS Investor Update - Dec 2024"
Slide 2: Revenue metrics (MRR, ARR, Net New MRR)
Slide 3: Unit economics (CAC, LTV, LTV/CAC Ratio, Churn Rate)

Use Dec 2024 data. Show trends vs Nov.
"""
    )
    if pptx_bytes:
        with open("SaaS_Investor_Update.pptx", "wb") as f:
            f.write(pptx_bytes)
        print("   Saved to SaaS_Investor_Update.pptx")
    else:
        print("   Failed to generate")

    # Deck 3
    print("3. Generating E-commerce Category Review...")
    pptx_bytes = generate_deck(
        ecommerce_data,
        """Create an E-commerce Review with exactly 2 slides:

Slide 1: Title slide - "E-commerce Category Review"
Slide 2: Top 3 categories by revenue with their YoY Growth %

Keep it simple - just highlight the top performers.
"""
    )
    if pptx_bytes:
        with open("Ecommerce_Category_Review.pptx", "wb") as f:
            f.write(pptx_bytes)
        print("   Saved to Ecommerce_Category_Review.pptx")
    else:
        print("   Failed to generate")

    print("\nDone! Open the .pptx files in PowerPoint or Keynote to view the decks.")


if __name__ == "__main__":
    main()
