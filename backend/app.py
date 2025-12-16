"""
Backend API for fetching LangSmith traces and rendering PPTX natively with Konva.
"""

from dotenv import load_dotenv
load_dotenv()

import os
import ast
import base64
from io import BytesIO
from typing import Optional
from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from langsmith import Client

# PPTX parsing and building
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

app = FastAPI(title="Slide Viewer API")

# CORS for React frontend
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:5173", "http://localhost:3000"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Initialize LangSmith client
ls_client = Client()


# ============================================================================
# PPTX PARSING - Extract slide data as JSON for Konva rendering
# ============================================================================

# Scale factor: PPTX uses inches, Konva uses pixels
# 96 DPI is standard for web
PIXELS_PER_INCH = 96

def emu_to_pixels(emu: int) -> float:
    """Convert EMUs (English Metric Units) to pixels."""
    inches = emu / 914400
    return inches * PIXELS_PER_INCH

def pixels_to_emu(pixels: float) -> int:
    """Convert pixels to EMUs."""
    inches = pixels / PIXELS_PER_INCH
    return int(inches * 914400)

def rgb_to_hex(rgb_color) -> Optional[str]:
    """Convert RGBColor or similar to hex string."""
    if rgb_color is None:
        return None
    try:
        if hasattr(rgb_color, 'rgb'):
            rgb_color = rgb_color.rgb
        if isinstance(rgb_color, RGBColor):
            return f"#{rgb_color.red:02x}{rgb_color.green:02x}{rgb_color.blue:02x}"
        return None
    except:
        return None

def hex_to_rgb(hex_color: str) -> tuple:
    """Convert hex color to RGB tuple."""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def get_background_color(slide) -> str:
    """Extract background color from slide."""
    try:
        background = slide.background
        fill = background.fill
        if fill.type is not None:
            fore_color = fill.fore_color
            if fore_color and fore_color.type is not None:
                color = rgb_to_hex(fore_color)
                if color:
                    return color
    except:
        pass
    return "#0f172a"  # Default dark blue

def get_text_alignment(alignment) -> str:
    """Convert pptx alignment to string."""
    if alignment is None:
        return "left"
    alignment_str = str(alignment).split(".")[-1].lower()
    if "center" in alignment_str:
        return "center"
    elif "right" in alignment_str:
        return "right"
    return "left"

def parse_shape(shape) -> Optional[dict]:
    """Parse a single shape into a JSON-serializable dict for Konva."""
    try:
        shape_data = {
            "id": str(shape.shape_id),
            "name": shape.name,
            "x": emu_to_pixels(shape.left),
            "y": emu_to_pixels(shape.top),
            "width": emu_to_pixels(shape.width),
            "height": emu_to_pixels(shape.height),
            "type": "rect",  # Default to rect
            "draggable": True,
        }
        
        # Handle text frames (textboxes)
        if shape.has_text_frame:
            text_content = shape.text_frame.text
            if text_content.strip():
                shape_data["type"] = "text"
                shape_data["text"] = text_content
                
                # Get text properties from first paragraph/run
                if shape.text_frame.paragraphs:
                    para = shape.text_frame.paragraphs[0]
                    shape_data["align"] = get_text_alignment(para.alignment)
                    
                    if para.runs:
                        run = para.runs[0]
                        if run.font.size:
                            shape_data["fontSize"] = run.font.size.pt
                        else:
                            shape_data["fontSize"] = 16
                        
                        color = rgb_to_hex(run.font.color)
                        shape_data["fill"] = color or "#ffffff"
                        shape_data["fontStyle"] = ""
                        if run.font.bold:
                            shape_data["fontStyle"] += "bold "
                        if run.font.italic:
                            shape_data["fontStyle"] += "italic"
                        shape_data["fontStyle"] = shape_data["fontStyle"].strip() or "normal"
                    else:
                        shape_data["fontSize"] = 16
                        shape_data["fill"] = "#ffffff"
                        shape_data["fontStyle"] = "normal"
        
        # Handle shapes (rectangles, etc.)
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            if shape_data["type"] != "text":
                shape_data["type"] = "rect"
            try:
                fill = shape.fill
                if fill.type is not None:
                    color = rgb_to_hex(fill.fore_color)
                    if color and shape_data["type"] == "rect":
                        shape_data["fill"] = color
            except:
                if shape_data["type"] == "rect":
                    shape_data["fill"] = "#1e293b"
        
        # Set default fill for rects without one
        if shape_data["type"] == "rect" and "fill" not in shape_data:
            shape_data["fill"] = "#1e293b"
        
        return shape_data
    except Exception as e:
        print(f"Error parsing shape: {e}")
        return None

def parse_pptx_bytes(pptx_bytes: bytes) -> dict:
    """
    Parse PPTX bytes into a JSON structure for Konva rendering.
    """
    try:
        prs = Presentation(BytesIO(pptx_bytes))
        
        result = {
            "width": emu_to_pixels(prs.slide_width),
            "height": emu_to_pixels(prs.slide_height),
            "slides": []
        }
        
        for idx, slide in enumerate(prs.slides):
            slide_data = {
                "id": f"slide-{idx}",
                "index": idx,
                "backgroundColor": get_background_color(slide),
                "shapes": []
            }
            
            for shape in slide.shapes:
                shape_data = parse_shape(shape)
                if shape_data:
                    slide_data["shapes"].append(shape_data)
            
            result["slides"].append(slide_data)
        
        return result
    except Exception as e:
        print(f"Error parsing PPTX: {e}")
        import traceback
        traceback.print_exc()
        return {"error": str(e), "slides": [], "width": 1280, "height": 720}


# ============================================================================
# PPTX BUILDING - Rebuild PPTX from edited JSON
# ============================================================================

def build_pptx_from_json(presentation_data: dict) -> bytes:
    """
    Build a PPTX file from JSON presentation data.
    """
    prs = Presentation()
    
    # Set slide dimensions
    width_px = presentation_data.get("width", 1280)
    height_px = presentation_data.get("height", 720)
    prs.slide_width = pixels_to_emu(width_px)
    prs.slide_height = pixels_to_emu(height_px)
    
    blank_layout = prs.slide_layouts[6]  # Blank slide
    
    for slide_data in presentation_data.get("slides", []):
        slide = prs.slides.add_slide(blank_layout)
        
        # Set background color
        bg_color = slide_data.get("backgroundColor", "#0f172a")
        r, g, b = hex_to_rgb(bg_color)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(r, g, b)
        
        # Add shapes
        for shape_data in slide_data.get("shapes", []):
            shape_type = shape_data.get("type", "rect")
            x = pixels_to_emu(shape_data.get("x", 0))
            y = pixels_to_emu(shape_data.get("y", 0))
            width = pixels_to_emu(shape_data.get("width", 100))
            height = pixels_to_emu(shape_data.get("height", 50))
            
            if shape_type == "text":
                # Add textbox
                textbox = slide.shapes.add_textbox(x, y, width, height)
                tf = textbox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = shape_data.get("text", "")
                
                # Text styling
                font_size = shape_data.get("fontSize", 16)
                p.font.size = Pt(font_size)
                
                fill_color = shape_data.get("fill", "#ffffff")
                r, g, b = hex_to_rgb(fill_color)
                p.font.color.rgb = RGBColor(r, g, b)
                
                font_style = shape_data.get("fontStyle", "normal")
                p.font.bold = "bold" in font_style
                p.font.italic = "italic" in font_style
                
                align = shape_data.get("align", "left")
                if align == "center":
                    p.alignment = PP_ALIGN.CENTER
                elif align == "right":
                    p.alignment = PP_ALIGN.RIGHT
                else:
                    p.alignment = PP_ALIGN.LEFT
                    
            elif shape_type == "rect":
                # Add rectangle
                shape = slide.shapes.add_shape(1, x, y, width, height)  # 1 = Rectangle
                fill_color = shape_data.get("fill", "#1e293b")
                r, g, b = hex_to_rgb(fill_color)
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(r, g, b)
                shape.line.fill.background()  # No border
    
    # Save to bytes
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


# ============================================================================
# PYDANTIC MODELS
# ============================================================================

class ShapeData(BaseModel):
    id: str
    name: Optional[str] = None
    x: float
    y: float
    width: float
    height: float
    type: str
    text: Optional[str] = None
    fontSize: Optional[float] = None
    fill: Optional[str] = None
    fontStyle: Optional[str] = None
    align: Optional[str] = None
    draggable: Optional[bool] = True

class SlideData(BaseModel):
    id: str
    index: int
    backgroundColor: str
    shapes: list[ShapeData]

class PresentationData(BaseModel):
    width: float
    height: float
    slides: list[SlideData]
    error: Optional[str] = None

class TraceSlide(BaseModel):
    trace_id: str
    trace_name: str
    created_at: str
    pptx_base64: Optional[str] = None
    presentation: Optional[PresentationData] = None
    error: Optional[str] = None

class TracesResponse(BaseModel):
    traces: list[TraceSlide]

class SavePresentationRequest(BaseModel):
    presentation: PresentationData

class SavePresentationResponse(BaseModel):
    pptx_base64: str
    success: bool


# ============================================================================
# API ENDPOINTS
# ============================================================================

def extract_pptx_from_trace(trace_id: str) -> Optional[bytes]:
    """Extract PPTX bytes from a trace's finalize_presentation tool call."""
    try:
        runs = list(ls_client.list_runs(
            project_name=os.getenv("LANGSMITH_PROJECT"),
            trace_id=trace_id,
        ))
        for run in runs:
            if run.name == "finalize_presentation" and run.outputs:
                output = run.outputs.get("output")
                if output:
                    content = output["content"]
                    pptx_bytes = ast.literal_eval(content)
                    print(f"Extracted {len(pptx_bytes)} bytes from trace")
                    return pptx_bytes
        print(f"No finalize_presentation output found in trace {trace_id}")
        return None
    except Exception as e:
        print(f"Error extracting PPTX from trace {trace_id}: {e}")
        import traceback
        traceback.print_exc()
        return None


@app.get("/api/traces", response_model=TracesResponse)
async def get_recent_traces():
    """Get the last 3 traces with their PPTX outputs, parsed for Konva rendering."""
    project_name = os.getenv("LANGSMITH_PROJECT", "default")
    try:
        root_runs = list(ls_client.list_runs(
            project_name=project_name,
            is_root=True,
            limit=3,
        ))
        
        result_traces = []
        for run in root_runs:
            trace_slide = TraceSlide(
                trace_id=str(run.trace_id),
                trace_name=run.name or "Unnamed",
                created_at=run.start_time.isoformat() if run.start_time else "",
            )
            
            pptx_bytes = extract_pptx_from_trace(str(run.trace_id))
            
            if pptx_bytes:
                # Parse PPTX for Konva rendering
                parsed = parse_pptx_bytes(pptx_bytes)
                if parsed.get("slides"):
                    trace_slide.presentation = PresentationData(**parsed)
                elif parsed.get("error"):
                    trace_slide.error = parsed["error"]
                
                # Keep base64 for download
                trace_slide.pptx_base64 = base64.b64encode(pptx_bytes).decode()
            else:
                trace_slide.error = "No PPTX output found in trace"
            
            result_traces.append(trace_slide)
        return TracesResponse(traces=result_traces)
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/save-presentation", response_model=SavePresentationResponse)
async def save_presentation(request: SavePresentationRequest):
    """Rebuild PPTX from edited presentation data and return as base64."""
    try:
        pptx_bytes = build_pptx_from_json(request.presentation.model_dump())
        pptx_base64 = base64.b64encode(pptx_bytes).decode()
        return SavePresentationResponse(pptx_base64=pptx_base64, success=True)
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/api/health")
async def health():
    return {
        "status": "ok",
        "langsmith_project": os.getenv("LANGSMITH_PROJECT", "default")
    }


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
